#!/usr/bin/env python3
import pandas as pd
import plotly.express as px
import os
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches

PPT_TEMPLATE_PATH = 'template.pptx'
DATA_FILE_PATH= 'testdata.csv'
GENERATED_PPT_PATH = 'new_presentation.pptx'


def create_slides(prs):
	title_slide_layout = prs.slide_layouts[0]
	prs.slides.add_slide(title_slide_layout)
	chart_and_table_slide_layout = prs.slide_layouts[5]
	prs.slides.add_slide(chart_and_table_slide_layout)	
	
	
def set_title_and_subtitle_to_slide(slide, title_text, subtitle_text=""):
	placeholders = {shape.name: shape.placeholder_format.idx for shape in slide.placeholders}
	try:
		title = slide.shapes.title
		title.text = title_text
		if 'Subtitle 2' in placeholders.keys() and subtitle_text:
			subtitle = slide.placeholders[1]
			subtitle.text = subtitle_text
	except Exception:
		print(f'Cannot set the title or subtitle to the given slide. \nValid placeholders: {placeholders}')


def create_and_manipulate_df():
	df = pd.read_csv(DATA_FILE_PATH)
	return df.groupby(['color'])['color'].count().reset_index(name="count")
	
	
def create_figure_and_save_to_file(df, figure_file_path):
	fig = px.bar(df, x="color", y="count", title="Color distribution")
	
	if not os.path.exists("images"):
		os.mkdir("images")
		
	fig.write_image(figure_file_path)
	
	
def put_picture_to_slide(slide, picture_path):
	placeholders = {shape.name: shape.placeholder_format.idx for shape in slide.placeholders}
	picture_placeholder_names = [name for name in placeholders.keys() if 'Picture' in name]
	try:
		picture_placeholder_index = placeholders[picture_placeholder_names[0]]
		picture_placeholder = slide.placeholders[picture_placeholder_index]
		picture = picture_placeholder.insert_picture(picture_path)
		picture.crop_top, picture.crop_bottom = 0, 0
		picture.crop_left, picture.crop_right = 0, 0
		
	except IndexError:
		print(f"Cannot put picture to the slide, there is no Picture placeholder in it. \nValid placeholders are: {placeholders}")


def put_table_to_slide(slide, df):
	placeholders = {shape.name: shape.placeholder_format.idx for shape in slide.placeholders}
	table_placeholder_names = [name for name in placeholders.keys() if 'Table' in name]
	try:
		table_placeholder_index = placeholders[table_placeholder_names[0]]
		table_placeholder = slide.placeholders[table_placeholder_index]
		if not df.empty:
			col_names = df.columns.tolist()
			rows = df.shape[0]+1
			cols = df.shape[1]
			shape = table_placeholder.insert_table(rows=rows, cols=cols)
			table = shape.table
			
			for i, col in enumerate(col_names):
				table.columns[i].width = Inches(2.0)
				table.cell(0, i).text = col
				for j in range(1, rows):
					table.cell(j, i).text = str(df[col].iloc[j-1])
	except IndexError:
		print(f"Cannot put table to the slide, there is no Table placeholder in it. \nValid placeholders are: {placeholders}")
	

prs_temp = Presentation(PPT_TEMPLATE_PATH)
create_slides(prs_temp)
set_title_and_subtitle_to_slide(prs_temp.slides[0], 'Main title', 'Main Subtitle')
df_distribution = create_and_manipulate_df()
create_figure_and_save_to_file(df_distribution, "images/fig1.png")
set_title_and_subtitle_to_slide(prs_temp.slides[1], "Table and chart title")
put_picture_to_slide(prs_temp.slides[1], "images/fig1.png")
put_table_to_slide(prs_temp.slides[1], df_distribution)

prs_temp.save(GENERATED_PPT_PATH)
print(f'Presentation generated successfully: {GENERATED_PPT_PATH}')
			